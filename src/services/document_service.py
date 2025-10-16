"""
Document processing service

Handles document extraction, processing, and chunking for various file formats
"""
import asyncio
import logging
import magic
import aiofiles
import email
import mailbox
from pathlib import Path
from typing import Dict, Any, Tuple, Optional, List
from fastapi import UploadFile

# Document processing libraries
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import xlrd
from bs4 import BeautifulSoup

from src.core.config import Settings
from src.models.schemas import DocumentType
from src.services.text_splitter import SimpleTextSplitter
from src.services.ocr_service import OCRService
from src.services.whatsapp_parser import WhatsAppParser
from src.services.llm_chat_parser import LLMChatParser
from src.services.email_threading_service import EmailThreadingService, EmailMessage
from src.services.document_type_handlers import (
    EmailHandler,
    ChatLogHandler,
)

logger = logging.getLogger(__name__)


class DocumentService:
    """
    Service for processing and extracting text from various document formats

    Supported formats:
    - PDF (text and scanned with OCR)
    - Office: Word (.docx, .doc), PowerPoint (.pptx), Excel (.xlsx, .xls)
    - Text: .txt, .md, code files
    - Email: .eml, .msg, .mbox (bulk archives)
    - Images: .png, .jpg, .jpeg, .tiff, .bmp (with OCR)
    - HTML/Webpages
    - WhatsApp exports
    """

    def __init__(self, settings: Settings):
        """
        Initialize document service

        Args:
            settings: Application settings
        """
        self.settings = settings
        self.text_splitter = SimpleTextSplitter(
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap
        )

        # Initialize OCR if enabled
        if settings.use_ocr:
            ocr_languages = settings.ocr_languages if isinstance(settings.ocr_languages, list) else settings.ocr_languages.split(',')
            self.ocr_service = OCRService(languages=ocr_languages)
        else:
            self.ocr_service = None

        # Initialize email threading service
        self.email_threading = EmailThreadingService()

        # Initialize document type handlers (only tested ones)
        self.email_handler = EmailHandler()
        self.chat_log_handler = ChatLogHandler()

        # Rate limiting: Max 5 concurrent document processing operations
        # Prevents OOM crashes from parallel enrichment calls
        self._processing_semaphore = asyncio.Semaphore(5)

    async def process_upload(
        self,
        file: UploadFile,
        process_ocr: bool = None
    ) -> Dict[str, Any]:
        """
        Process uploaded file

        Args:
            file: Uploaded file from FastAPI
            process_ocr: Force OCR processing (overrides settings)

        Returns:
            Dictionary with extracted content and metadata
        """
        # Save uploaded file temporarily
        temp_path = Path(self.settings.temp_path) / file.filename
        temp_path.parent.mkdir(parents=True, exist_ok=True)

        try:
            async with aiofiles.open(temp_path, 'wb') as f:
                content = await file.read()
                await f.write(content)

            # Rate-limited processing: max 5 concurrent operations
            async with self._processing_semaphore:
                # Process the file
                text, doc_type, metadata = await self.extract_text_from_file(
                    file_path=temp_path,
                    process_ocr=process_ocr if process_ocr is not None else self.settings.use_ocr
                )

                # Apply document type-specific handler for preprocessing
                text, metadata = self._apply_document_type_handler(text, doc_type, metadata, temp_path)

                # Add file metadata
                metadata.update({
                    "original_filename": file.filename,
                    "content_type": file.content_type,
                    "file_size_bytes": len(content)
                })

            return {
                "text": text,
                "document_type": doc_type,
                "metadata": metadata,
                "filename": file.filename
            }

        finally:
            # Cleanup temp file
            if temp_path.exists():
                temp_path.unlink()

    async def extract_text_from_file(
        self,
        file_path: str | Path,
        process_ocr: bool = False
    ) -> Tuple[str, DocumentType, Dict[str, Any]]:
        """
        Extract text from various file formats

        Args:
            file_path: Path to file
            process_ocr: Enable OCR for images and scanned PDFs

        Returns:
            Tuple of (extracted_text, document_type, metadata)

        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If file is too large or unsupported format
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        # Check file size
        file_size_mb = file_path.stat().st_size / (1024 * 1024)
        if file_size_mb > self.settings.max_file_size_mb:
            raise ValueError(
                f"File too large: {file_size_mb:.1f}MB (max: {self.settings.max_file_size_mb}MB)"
            )

        # Detect file type
        mime_type = magic.from_file(str(file_path), mime=True)
        file_extension = file_path.suffix.lower()

        logger.info(f"Processing file: {file_path.name}, MIME: {mime_type}, Ext: {file_extension}")

        metadata = {
            "file_extension": file_extension,
            "mime_type": mime_type,
            "file_size_mb": round(file_size_mb, 2)
        }

        # Route to appropriate processor
        if mime_type == "application/pdf" or file_extension == ".pdf":
            text, doc_type = await self._process_pdf(file_path, process_ocr)
            return text, doc_type, metadata

        elif file_extension in ['.docx', '.doc']:
            text, doc_metadata = await self._process_word_document(file_path)
            metadata.update(doc_metadata)
            return text, DocumentType.office, metadata

        elif file_extension in ['.pptx', '.ppt']:
            text = await self._process_powerpoint(file_path)
            return text, DocumentType.office, metadata

        elif file_extension in ['.xlsx', '.xls']:
            text = await self._process_excel(file_path)
            return text, DocumentType.office, metadata

        elif mime_type.startswith("image/") or file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
            if process_ocr and self.ocr_service:
                text = self.ocr_service.extract_text_from_image(file_path)
                return text, DocumentType.scanned, metadata
            else:
                return f"Image file: {file_path.name}", DocumentType.image, metadata

        elif file_extension in ['.eml', '.msg']:
            text, email_metadata = await self._process_email(file_path)
            logger.info(f"📧 RECEIVED email_metadata with keys: {list(email_metadata.keys())}, thread_id={email_metadata.get('thread_id', 'NONE')}")
            metadata.update(email_metadata)
            logger.info(f"📧 AFTER UPDATE, metadata has keys: {list(metadata.keys())}, thread_id={metadata.get('thread_id', 'NONE')}")
            return text, DocumentType.email, metadata

        elif file_extension == '.mbox':
            text = await self._process_mbox(file_path)
            return text, DocumentType.email, metadata

        elif mime_type.startswith("text/html") or file_extension in ['.html', '.htm']:
            text = await self._process_html(file_path)
            return text, DocumentType.webpage, metadata

        elif (mime_type.startswith("text/") or
              file_extension in ['.txt', '.md', '.py', '.js', '.json', '.yaml', '.yml', '.xml', '.css']):
            async with aiofiles.open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = await f.read()

            # Check for special text formats
            if WhatsAppParser.is_whatsapp_export(content):
                # Parse and thread WhatsApp messages
                messages, summary, wa_metadata = WhatsAppParser.parse_whatsapp_export(content)

                if messages:
                    # Group into conversation threads (4-hour gaps)
                    threads = WhatsAppParser.group_into_threads(messages, time_gap_hours=4)
                    logger.info(f"WhatsApp: {len(messages)} messages → {len(threads)} threads")

                    # Format as threaded conversation
                    result = f"WhatsApp Export: {file_path.name}\n"
                    result += f"{summary}\n\n"

                    for thread_idx, thread in enumerate(threads, 1):
                        result += WhatsAppParser.format_thread_as_text(thread, thread_idx)
                        result += "\n"

                    # Merge metadata
                    metadata.update(wa_metadata)

                    return result, DocumentType.whatsapp, metadata
                else:
                    # Fallback if parsing fails
                    return content, DocumentType.whatsapp, metadata

            elif LLMChatParser.is_llm_export(content):
                # Parse LLM chat export (ChatGPT, Claude, markdown, etc.)
                messages, summary, llm_metadata = LLMChatParser.parse_llm_export(content)

                if messages:
                    logger.info(f"LLM Chat: {len(messages)} messages from {llm_metadata.get('export_type', 'unknown')} export")

                    # Format as markdown conversation
                    conv_title = llm_metadata.get("conversation_title") or llm_metadata.get("conversation_titles", ["LLM Conversation"])[0]
                    result = LLMChatParser.format_as_markdown(messages, conv_title)

                    # Add summary header
                    result = f"{summary}\n\n{'='*80}\n\n{result}"

                    # Merge metadata
                    metadata.update(llm_metadata)

                    # Return as LLM chat document type
                    return result, DocumentType.llm_chat, metadata
                else:
                    # For markdown format, content is already formatted - just return with llm_chat type
                    # The ChunkingService will handle parsing into turns
                    if llm_metadata.get('export_type') == 'markdown':
                        logger.info(f"LLM Chat (markdown): {llm_metadata.get('message_count', 0)} messages detected")
                        metadata.update(llm_metadata)
                        return content, DocumentType.llm_chat, metadata
                    else:
                        # Fallback if parsing fails for other formats
                        return content, DocumentType.text, metadata

            elif file_extension in ['.py', '.js', '.java', '.cpp', '.c', '.cs', '.php', '.go', '.rs']:
                return content, DocumentType.code, metadata
            else:
                return content, DocumentType.text, metadata

        else:
            raise ValueError(f"Unsupported file type: {mime_type} ({file_extension})")

    async def _process_pdf(self, file_path: Path, process_ocr: bool) -> Tuple[str, DocumentType]:
        """
        Process PDF files with optional OCR fallback

        Args:
            file_path: Path to PDF file
            process_ocr: Enable OCR if text extraction fails

        Returns:
            Tuple of (extracted_text, document_type)
        """
        try:
            # Try text extraction first
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)

                logger.info(f"Extracting text from {num_pages} PDF pages")

                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text += page_text + "\n"

            # If no text extracted or OCR requested, use OCR
            if (not text.strip() or process_ocr) and self.ocr_service and self.ocr_service.is_available():
                logger.info(f"Using OCR for PDF: {file_path.name}")
                ocr_text = self.ocr_service.extract_text_from_pdf_images(file_path)
                text = ocr_text if ocr_text.strip() else text
                return text, DocumentType.scanned

            return text, DocumentType.pdf

        except Exception as e:
            logger.error(f"PDF processing failed for {file_path}: {e}")

            # Fallback to OCR if available
            if self.ocr_service and self.ocr_service.is_available():
                logger.info("Falling back to OCR")
                return self.ocr_service.extract_text_from_pdf_images(file_path), DocumentType.scanned
            else:
                raise

    async def _process_word_document(self, file_path: Path) -> Tuple[str, Dict]:
        """
        Process Word documents (.docx, .doc)

        Args:
            file_path: Path to Word file

        Returns:
            Tuple of (extracted_text, metadata_dict)
        """
        try:
            doc = DocxDocument(str(file_path))

            # Extract paragraphs
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])

            # Extract metadata
            metadata = {}
            if doc.core_properties.title:
                metadata["title"] = doc.core_properties.title
            if doc.core_properties.author:
                metadata["author"] = doc.core_properties.author
            if doc.core_properties.created:
                metadata["created"] = doc.core_properties.created.isoformat()

            logger.info(f"Successfully extracted text from Word document: {file_path.name}")
            return text, metadata

        except Exception as e:
            logger.error(f"Word document processing failed for {file_path}: {e}")
            return f"Failed to process Word document: {file_path.name}", {}

    async def _process_powerpoint(self, file_path: Path) -> str:
        """Process PowerPoint presentations"""
        try:
            prs = Presentation(str(file_path))
            text = ""

            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slide {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

            return text

        except Exception as e:
            logger.error(f"PowerPoint processing failed for {file_path}: {e}")
            return f"Failed to process PowerPoint: {file_path.name}"

    async def _process_excel(self, file_path: Path) -> str:
        """Process Excel spreadsheets (.xlsx, .xls)"""
        try:
            text = ""

            if file_path.suffix.lower() == '.xlsx':
                workbook = openpyxl.load_workbook(str(file_path), data_only=True)

                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text += f"\n--- Sheet: {sheet_name} ---\n"

                    for row in sheet.iter_rows(values_only=True):
                        row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                        if row_text.strip():
                            text += row_text + "\n"
            else:  # .xls
                workbook = xlrd.open_workbook(str(file_path))

                for sheet in workbook.sheets():
                    text += f"\n--- Sheet: {sheet.name} ---\n"

                    for row_idx in range(sheet.nrows):
                        row = sheet.row_values(row_idx)
                        row_text = "\t".join([str(cell) for cell in row if cell])
                        if row_text.strip():
                            text += row_text + "\n"

            return text

        except Exception as e:
            logger.error(f"Excel processing failed for {file_path}: {e}")
            return f"Failed to process Excel file: {file_path.name}"

    async def _process_email(self, file_path: Path) -> Tuple[str, Dict]:
        """
        Process email files (.eml, .msg) with proper charset handling

        Extracts:
        - Email headers (From, To, Subject, Date)
        - Email body (text/plain and text/html)
        - Attachments (saved and referenced for separate processing)
        - Threading metadata (Message-ID, In-Reply-To, References)

        Returns:
            Tuple of (email text content, metadata dict with created_date + threading fields)
        """
        try:
            import tempfile
            import os
            import email.utils
            from datetime import date

            with open(file_path, 'rb') as f:
                msg = email.message_from_bytes(f.read())

            # Parse email date for metadata
            metadata = {}
            date_str = msg.get('Date')
            if date_str:
                try:
                    # email.utils.parsedate_to_datetime handles RFC 2822 date format
                    email_datetime = email.utils.parsedate_to_datetime(date_str)
                    # Convert to date string for metadata
                    metadata['created_date'] = email_datetime.date().isoformat()
                    logger.debug(f"Parsed email date: {date_str} → {metadata['created_date']}")
                except Exception as e:
                    logger.warning(f"Failed to parse email date '{date_str}': {e}")

            # Extract threading metadata
            metadata['message_id'] = msg.get('Message-ID', f"<generated-{file_path.stem}>")
            metadata['in_reply_to'] = msg.get('In-Reply-To', '')
            metadata['references'] = msg.get('References', '')

            # Extract thread topic/index if available
            thread_topic = msg.get('Thread-Topic', '')
            if thread_topic:
                metadata['thread_topic'] = thread_topic
            thread_index = msg.get('Thread-Index', '')
            if thread_index:
                metadata['thread_index'] = thread_index

            # Extract sender and recipients for threading
            metadata['sender'] = msg.get('From', 'Unknown')
            metadata['recipients'] = ', '.join(msg.get_all('To', []) or ['Unknown'])

            # Subject for threading
            subject = msg.get('Subject', '(No Subject)')
            metadata['subject'] = subject

            # Generate thread_id from normalized subject
            normalized_subject = self.email_threading.normalize_subject(subject)
            import hashlib
            metadata['thread_id'] = hashlib.md5(normalized_subject.encode()).hexdigest()[:12]

            logger.info(f"📧 Threading: thread_id={metadata['thread_id']}, message_id={metadata['message_id'][:50]}...")

            # Extract headers with proper decoding
            def decode_header_value(header_value):
                """Decode email header, handling bytes and encoded-words"""
                if not header_value:
                    return 'Unknown'
                decoded_parts = email.header.decode_header(header_value)
                result = []
                for part, charset in decoded_parts:
                    if isinstance(part, bytes):
                        result.append(part.decode(charset or 'utf-8', errors='replace'))
                    else:
                        result.append(str(part))
                return ''.join(result)

            text = f"From: {decode_header_value(msg.get('From'))}\n"
            text += f"To: {decode_header_value(msg.get('To'))}\n"
            text += f"Subject: {decode_header_value(msg.get('Subject'))}\n"
            text += f"Date: {msg.get('Date', 'Unknown')}\n\n"

            # Track attachments - save to persistent storage
            attachments = []
            # Save to processed directory (so they persist)
            from datetime import datetime
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            attachment_dir = Path(self.settings.processed_path) / "email_attachments" / f"{timestamp}_{file_path.stem}"
            attachment_dir.mkdir(parents=True, exist_ok=True)
            logger.info(f"📎 Saving email attachments to: {attachment_dir}")

            # Extract body and attachments
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    content_disposition = str(part.get("Content-Disposition", ""))

                    # Extract text body
                    if content_type == "text/plain" and "attachment" not in content_disposition:
                        payload = part.get_payload(decode=True)
                        if payload:
                            charset = part.get_content_charset() or 'utf-8'
                            try:
                                text += payload.decode(charset, errors='replace')
                            except (LookupError, UnicodeDecodeError):
                                for fallback_charset in ['utf-8', 'iso-8859-1', 'windows-1252']:
                                    try:
                                        text += payload.decode(fallback_charset, errors='replace')
                                        break
                                    except:
                                        continue

                    # Extract HTML body as fallback (convert to text)
                    elif content_type == "text/html" and "attachment" not in content_disposition:
                        payload = part.get_payload(decode=True)
                        if payload and not text.strip():  # Only if no plain text
                            charset = part.get_content_charset() or 'utf-8'
                            try:
                                html_content = payload.decode(charset, errors='replace')
                                # Simple HTML to text conversion
                                soup = BeautifulSoup(html_content, 'html.parser')
                                text += soup.get_text(separator='\n', strip=True)
                            except:
                                pass

                    # Save attachments
                    elif "attachment" in content_disposition or part.get_filename():
                        filename = part.get_filename()
                        if filename:
                            # Decode filename if encoded
                            if isinstance(filename, str):
                                decoded_filename = email.header.decode_header(filename)
                                filename_parts = []
                                for fname_part, fcharset in decoded_filename:
                                    if isinstance(fname_part, bytes):
                                        filename_parts.append(fname_part.decode(fcharset or 'utf-8', errors='replace'))
                                    else:
                                        filename_parts.append(fname_part)
                                filename = ''.join(filename_parts)

                            # Save attachment
                            attachment_path = attachment_dir / filename
                            payload = part.get_payload(decode=True)
                            if payload:
                                with open(attachment_path, 'wb') as att_file:
                                    att_file.write(payload)
                                attachments.append(str(attachment_path))
                                logger.info(f"📎 Saved email attachment: {filename} ({len(payload)} bytes)")
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    charset = msg.get_content_charset() or 'utf-8'
                    try:
                        text += payload.decode(charset, errors='replace')
                    except (LookupError, UnicodeDecodeError):
                        for fallback_charset in ['utf-8', 'iso-8859-1', 'windows-1252']:
                            try:
                                text += payload.decode(fallback_charset, errors='replace')
                                break
                            except:
                                continue

            # Add attachment references to email text
            if attachments:
                text += "\n\n--- Attachments ---\n"
                for att_path in attachments:
                    text += f"📎 {Path(att_path).name} (saved at: {att_path})\n"

                # Store attachment paths in metadata for batch processing
                metadata['attachment_paths'] = attachments
                metadata['attachment_count'] = len(attachments)
                metadata['has_attachments'] = True

                logger.info(f"📎 Email has {len(attachments)} attachments for processing")
            else:
                metadata['has_attachments'] = False
                metadata['attachment_count'] = 0

            # Preprocess email content using EmailHandler
            # This removes reply chains, forwarding headers, signatures, etc.
            original_length = len(text)
            metadata['original_length'] = original_length  # For handler to calculate removal ratio

            text = self.email_handler.preprocess(text, metadata)

            # Extract additional email-specific metadata
            handler_metadata = self.email_handler.extract_metadata(text, metadata)
            metadata.update(handler_metadata)

            logger.info(
                f"📧 EMAIL METADATA BEING RETURNED: thread_id={metadata.get('thread_id', 'NONE')}, "
                f"cleaned: {original_length}→{len(text)} chars, keys={list(metadata.keys())}"
            )
            return text, metadata

        except Exception as e:
            logger.error(f"Email processing failed for {file_path}: {e}")
            return f"Failed to process email: {file_path.name}", {}

    async def _process_mbox(self, file_path: Path) -> str:
        """
        Process mbox archive file containing multiple emails
        Groups emails into conversation threads for better context

        Args:
            file_path: Path to .mbox file

        Returns:
            Text organized by conversation threads
        """
        try:
            from datetime import datetime
            import email.utils

            mbox = mailbox.mbox(str(file_path))
            email_messages = []

            logger.info(f"Processing mbox archive: {file_path.name}")

            # Parse all emails into EmailMessage objects
            for idx, message in enumerate(mbox, 1):
                try:
                    # Extract message ID
                    message_id = message.get('Message-ID', f"<generated-{idx}>")

                    # Extract headers
                    subject = message.get('Subject', '(No Subject)')
                    sender = message.get('From', 'Unknown')
                    recipients = message.get_all('To', [])
                    if isinstance(recipients, str):
                        recipients = [recipients]

                    # Parse date
                    date_str = message.get('Date')
                    try:
                        date = email.utils.parsedate_to_datetime(date_str) if date_str else datetime.now()
                    except Exception:
                        date = datetime.now()

                    # Extract In-Reply-To and References
                    in_reply_to = message.get('In-Reply-To')
                    references = message.get_all('References', [])
                    if isinstance(references, str):
                        references = [references]

                    # Extract body
                    body = ""
                    if message.is_multipart():
                        for part in message.walk():
                            if part.get_content_type() == "text/plain":
                                payload = part.get_payload(decode=True)
                                if payload:
                                    body += payload.decode('utf-8', errors='ignore')
                    else:
                        payload = message.get_payload(decode=True)
                        if payload:
                            body = payload.decode('utf-8', errors='ignore')

                    # Create EmailMessage object
                    email_msg = EmailMessage(
                        message_id=message_id,
                        subject=subject,
                        sender=sender,
                        recipients=recipients,
                        date=date,
                        body=body,
                        in_reply_to=in_reply_to,
                        references=references
                    )
                    email_messages.append(email_msg)

                except Exception as e:
                    logger.warning(f"Failed to parse email {idx} in mbox: {e}")
                    continue

            logger.info(f"Parsed {len(email_messages)} emails from {file_path.name}")

            # Group into threads
            threads = self.email_threading.build_threads(email_messages)
            logger.info(f"Grouped into {len(threads)} conversation threads")

            # Format output by thread
            result = f"MBOX Archive: {file_path.name}\n"
            result += f"Total Emails: {len(email_messages)}\n"
            result += f"Conversation Threads: {len(threads)}\n\n"

            for thread_idx, thread in enumerate(threads, 1):
                result += f"\n{'='*80}\n"
                result += f"THREAD {thread_idx}: {thread.subject}\n"
                result += f"{'='*80}\n"
                result += f"Messages: {thread.message_count}\n"
                result += f"Participants: {', '.join(sorted(thread.participants))}\n"
                result += f"Date Range: {thread.start_date.strftime('%Y-%m-%d')} to {thread.end_date.strftime('%Y-%m-%d')}\n"
                result += f"{'-'*80}\n\n"

                for msg_idx, msg in enumerate(thread.messages, 1):
                    result += f"[Message {msg_idx}/{thread.message_count}]\n"
                    result += f"From: {msg.sender}\n"
                    result += f"To: {', '.join(msg.recipients)}\n"
                    result += f"Date: {msg.date.strftime('%Y-%m-%d %H:%M')}\n"
                    result += f"{'-'*40}\n"
                    result += f"{msg.body}\n\n"

            return result

        except Exception as e:
            logger.error(f"MBOX processing failed for {file_path}: {e}")
            return f"Failed to process mbox archive: {file_path.name}"

    async def _process_html(self, file_path: Path) -> str:
        """Process HTML files and extract text"""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = await f.read()

            soup = BeautifulSoup(html_content, 'html.parser')

            # Remove script and style elements
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()

            # Get text
            text = soup.get_text(separator='\n', strip=True)

            # Clean up extra whitespace
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            text = '\n'.join(lines)

            return text

        except Exception as e:
            logger.error(f"HTML processing failed for {file_path}: {e}")
            return f"Failed to process HTML: {file_path.name}"

    def _apply_document_type_handler(
        self,
        text: str,
        doc_type: DocumentType,
        metadata: Dict[str, Any],
        file_path: Optional[Path] = None
    ) -> Tuple[str, Dict[str, Any]]:
        """
        Apply type-specific document handler for preprocessing and metadata extraction.

        Currently supports only TESTED handlers:
        - Email handler (German emails validated)
        - Chat log handler (ChatGPT format validated)

        Args:
            text: Extracted text
            doc_type: Detected document type
            metadata: Existing metadata
            file_path: Optional file path for additional context

        Returns:
            Tuple of (preprocessed_text, updated_metadata)
        """
        original_length = len(text)
        handler_used = None

        # Email handler (already applied in _process_email, but this is for consistency)
        if doc_type == DocumentType.email:
            # Email preprocessing already done in _process_email
            handler_used = "email"

        # Chat log handler (ChatGPT exports)
        elif doc_type == DocumentType.llm_chat:
            metadata['original_length'] = original_length
            text = self.chat_log_handler.preprocess(text, metadata)
            handler_metadata = self.chat_log_handler.extract_metadata(text, metadata)
            metadata.update(handler_metadata)
            handler_used = "chat_log"

        # Note: Invoice, scanned, and manual handlers moved to src/services/experimental/
        # These are untested and should not be used in production

        if handler_used:
            cleaned_length = len(text)
            retention_pct = (cleaned_length / original_length * 100) if original_length > 0 else 0
            logger.info(f"📋 Applied {handler_used} handler: {original_length} → {cleaned_length} chars ({retention_pct:.1f}% retained)")
            metadata['handler_applied'] = handler_used

        return text, metadata

    def chunk_text(self, text: str) -> List[str]:
        """
        Split text into chunks for processing

        Args:
            text: Text to split

        Returns:
            List of text chunks
        """
        return self.text_splitter.split_text(text)

    def clean_content(self, content: str) -> str:
        """
        Clean and normalize text content

        Args:
            content: Raw text content

        Returns:
            Cleaned text
        """
        import re

        # Remove null bytes
        content = content.replace('\x00', '')

        # Normalize whitespace
        lines = content.split('\n')
        cleaned_lines = []
        previous_was_empty = False

        for line in lines:
            # Strip trailing whitespace
            line = line.rstrip()

            # Collapse multiple spaces to single space
            line = re.sub(r' +', ' ', line)

            # Handle empty lines (preserve one for paragraph separation)
            if not line.strip():
                if not previous_was_empty:
                    cleaned_lines.append('')
                    previous_was_empty = True
            else:
                cleaned_lines.append(line)
                previous_was_empty = False

        # Remove trailing empty lines
        while cleaned_lines and not cleaned_lines[-1].strip():
            cleaned_lines.pop()

        return '\n'.join(cleaned_lines)
